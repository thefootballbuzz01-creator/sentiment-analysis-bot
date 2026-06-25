#!/usr/bin/env python3
"""Builds the sentiment-analysis presentation as a STORY, with python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ---- palette (ocean / teal) ----
NAVY  = "0B2942"
TEAL  = "0D9488"
CYAN  = "22D3EE"
ICE   = "CFE9EF"
LIGHT = "F2F7F9"
CARD  = "FFFFFF"
INK   = "11242F"
MUTED = "5B7383"
WHITE = "FFFFFF"
POS   = "14B8A6"
NEU   = "38BDF8"
NEG   = "FB7185"
YT_C  = "F87171"   # YouTube identity (red)
GP_C  = "34D399"   # Google Play identity (green)

HEAD = "Cambria"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def rgb(h): return RGBColor.from_string(h)


def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(bg)
    return s


def soft_shadow(shape):
    spPr = shape._element.spPr
    e = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
                          {'blurRad': '80000', 'dist': '25000',
                           'dir': '5400000', 'rotWithShape': '0'})
    c = spPr.makeelement(qn('a:srgbClr'), {'val': '0B2942'})
    a = spPr.makeelement(qn('a:alpha'), {'val': '15000'})
    c.append(a); sh.append(c); e.append(sh); spPr.append(e)


def text(s, txt, x, y, w, h, size=18, color=INK, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
         spacing=None, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(txt if isinstance(txt, list) else [txt]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.name = font; f.color.rgb = rgb(color)
        if spacing is not None:
            r._r.get_or_add_rPr().set('spc', str(int(spacing * 100)))
    return tb


def card(s, x, y, w, h, fill=CARD, radius=0.09, shadow=True, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line:
        sp.line.color.rgb = rgb(line); sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if shadow:
        soft_shadow(sp)
    else:
        sp.shadow.inherit = False
    return sp


def circle_num(s, n, x, y, d=0.62, fill=TEAL, tc=WHITE, size=20):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    sp.line.fill.background(); sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(size); r.font.bold = True; r.font.name = HEAD
    r.font.color.rgb = rgb(tc)
    return sp


def dot(s, x, y, d, color):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(color)
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def bar(s, x, y, wmax, frac, color):
    tr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                            Inches(wmax), Inches(0.16))
    tr.fill.solid(); tr.fill.fore_color.rgb = rgb("E2E8F0")
    tr.line.fill.background(); tr.shadow.inherit = False
    try: tr.adjustments[0] = 0.5
    except Exception: pass
    fl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                            Inches(max(0.12, wmax * frac)), Inches(0.16))
    fl.fill.solid(); fl.fill.fore_color.rgb = rgb(color)
    fl.line.fill.background(); fl.shadow.inherit = False
    try: fl.adjustments[0] = 0.5
    except Exception: pass


def quote_pill(s, x, y, w, text_str, color, label):
    card(s, x, y, w, 1.15, fill=NAVY, shadow=True)
    text(s, label, x + 0.3, y + 0.18, w - 0.6, 0.3, 12, color, True, BODY)
    text(s, text_str, x + 0.3, y + 0.5, w - 0.6, 0.55, 15, WHITE, False, BODY,
         italic=True)


# ============================================================ 1. HOOK
s = slide(NAVY)
text(s, "A DATA STORY", 0.95, 0.85, 10, 0.4, 15, CYAN, True, HEAD, spacing=3)
text(s, "What are customers\nREALLY saying?".split("\n"), 0.9, 1.4, 11.6, 2.0,
     50, WHITE, True, HEAD, line_spacing=1.0)
text(s, "I taught a computer to read 698 real comments about Argos — for free. "
        "What it heard surprised me.", 0.95, 3.6, 11.0, 0.8, 20, ICE, False,
     BODY, italic=True)
quote_pill(s, 0.95, 4.7, 5.6, "“Easy to use — find what you want straight away.”",
           "6EE7B7", "ONE CUSTOMER")
quote_pill(s, 6.75, 4.7, 5.6, "“Worst app in the business. Items just disappear.”",
           "FCA5A5", "ANOTHER, SAME DAY")
text(s, "Presented by  ______________      •      "
        "YouTube + Google Play + AI", 0.95, 6.55, 11.5, 0.5, 14, "9FC6D3",
     False, BODY)
s.notes_slide.notes_text_frame.text = (
    "Open on the hook — don't explain the tech yet. Read the two quotes out loud: "
    "the SAME product, love and fury on the same day. The question 'what are "
    "customers really saying?' is the thread for the whole talk.")

# ============================================================ 2. THE PROBLEM
s = slide(LIGHT)
text(s, "The truth was hiding in plain sight", 0.7, 0.7, 12, 0.9, 38, INK,
     True, HEAD)
text(s, "Every day, customers pour their feelings into comments and reviews — "
        "praise, frustration, fury. But it's scattered across hundreds of "
        "videos and app reviews. Nobody can read it all. So nobody really "
        "knows what customers think.", 0.7, 1.7, 11.9, 1.4, 19, MUTED, False,
     BODY, line_spacing=1.2)
big = card(s, 0.7, 3.5, 11.93, 2.7, fill=NAVY, shadow=True)
text(s, "698", 0.7, 3.95, 5.9, 1.4, 90, CYAN, True, HEAD, align=PP_ALIGN.CENTER)
text(s, "real voices about Argos", 0.7, 5.35, 5.9, 0.5, 18, ICE, False, BODY,
     align=PP_ALIGN.CENTER)
text(s, "0", 6.7, 3.95, 5.9, 1.4, 90, NEG, True, HEAD, align=PP_ALIGN.CENTER)
text(s, "people reading them all", 6.7, 5.35, 5.9, 0.5, 18, ICE, False, BODY,
     align=PP_ALIGN.CENTER)
s.notes_slide.notes_text_frame.text = (
    "Build the tension. The feedback EXISTS — it's just unreadable at human "
    "scale. 698 voices, nobody listening. That gap is the problem the project "
    "solves. Pause on the 698-vs-0.")

# ============================================================ 3. THE QUEST
s = slide(NAVY)
text(s, "So I built something to listen", 0.7, 0.85, 12, 0.9, 38, WHITE, True, HEAD)
text(s, "No budget. No team. No paid services. Just code — a bot that finds the "
        "comments, reads the feeling behind every single one, and tells the "
        "story.", 0.7, 1.85, 11.9, 1.1, 20, ICE, False, BODY, line_spacing=1.2)
steps = [("1", "Find", "real comments"), ("2", "Read", "the sentiment"),
         ("3", "Store", "every voice"), ("4", "Show", "the live story"),
         ("5", "Repeat", "every hour")]
cw, gap, sx = 2.19, 0.26, 0.7
for i, (n, hd, sub) in enumerate(steps):
    x = sx + i * (cw + gap)
    card(s, x, 3.3, cw, 2.3, fill="10324C", shadow=False, line="1C4A5E")
    circle_num(s, n, x + cw/2 - 0.31, 3.62, 0.62, CYAN, NAVY)
    text(s, hd, x + 0.15, 4.45, cw - 0.3, 0.45, 19, WHITE, True, HEAD,
         align=PP_ALIGN.CENTER)
    text(s, sub, x + 0.15, 4.95, cw - 0.3, 0.4, 13, "9FC6D3", False, BODY,
         align=PP_ALIGN.CENTER)
text(s, "It still runs right now — on its own, every hour, for £0.",
     0.7, 5.95, 12, 0.5, 17, CYAN, True, BODY, italic=True)
s.notes_slide.notes_text_frame.text = (
    "This is the 'I set out on a quest' beat. Emphasise the constraints — no "
    "money, no team — because the payoff is that it still works. Walk the five "
    "verbs quickly; the detail comes next.")

# ============================================================ 4. HOW IT FEELS
s = slide(LIGHT)
text(s, "How do you teach a computer to feel?", 0.7, 0.6, 12, 0.9, 36, INK,
     True, HEAD)
text(s, "A tool called VADER scores every comment from −1 (furious) to +1 "
        "(delighted). It even understands sarcasm, CAPITALS and emoji.",
     0.7, 1.55, 11.9, 0.9, 18, MUTED, False, BODY, line_spacing=1.15)
ex = [("“Love this — fast delivery!”", "+0.84", "Delighted", POS),
      ("“It's okay, nothing special.”", "0.00", "Neutral", NEU),
      ("“Items just disappear. Useless.”", "−0.71", "Furious", NEG)]
y = 2.75
for q, sc, lab, col in ex:
    card(s, 0.7, y, 8.4, 0.94, fill=CARD)
    text(s, q, 1.05, y, 6.0, 0.94, 17, INK, False, BODY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, sc, 7.0, y, 1.7, 0.94, 22, col, True, HEAD, anchor=MSO_ANCHOR.MIDDLE,
         align=PP_ALIGN.CENTER)
    pill = card(s, 9.3, y + 0.22, 3.33, 0.5, fill=col, radius=0.5, shadow=False)
    text(s, lab, 9.3, y + 0.22, 3.33, 0.5, 15, WHITE, True, BODY,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    y += 1.14
text(s, "No training, no AI bill — it runs instantly, and you can see exactly "
        "why it chose each label.", 0.7, 6.4, 11.9, 0.5, 15, MUTED, True, BODY)
s.notes_slide.notes_text_frame.text = (
    "Read the three examples in three different tones of voice. WHY VADER: "
    "(1) it's FREE — no training, no API, no internet — which fits the zero-cost "
    "project; (2) it was built specifically for short SOCIAL-MEDIA text, so it "
    "handles the slang, CAPITALS, punctuation and emoji in comments and reviews "
    "that trip up general tools; (3) it's EXPLAINABLE — rule + dictionary based, "
    "so you can show exactly which words drove each score (a black-box AI can't); "
    "(4) it's INSTANT and lightweight, so it scales to hundreds of comments and "
    "runs hourly for free. Honest limit: it can miss sarcasm and non-English "
    "text. Why not alternatives? A trained ML model or an LLM scoring every "
    "comment would be more nuanced but cost money/compute and aren't "
    "transparent — VADER is the sweet spot for this exact kind of text.")

# ============================================================ 5. PLOT TWIST
s = slide(LIGHT)
text(s, "Plot twist: the data lied to me", 0.7, 0.7, 12, 0.9, 38, INK, True, HEAD)
text(s, "My first results were full of… wristwatches and aftershave.",
     0.7, 1.7, 11.9, 0.6, 21, NEG, True, BODY, italic=True)
text(s, "It turns out “Argos” is also a watch brand and a perfume line. Nearly "
        "half my “data” wasn't about the shop at all. A good analyst doubts the "
        "evidence — so I cleaned it: I retrained the search on retail words and "
        "filtered the noise out. Only then did the real story appear.",
     0.7, 2.55, 11.9, 1.6, 19, MUTED, False, BODY, line_spacing=1.25)
hl = card(s, 0.7, 4.5, 11.93, 1.5, fill=NAVY, shadow=True)
text(s, "Neutral comments dropped by 60% once the noise was gone — proof the "
        "clean-up worked.", 1.0, 4.5, 11.3, 1.5, 20, WHITE, True, BODY,
     anchor=MSO_ANCHOR.MIDDLE, italic=True)
s.notes_slide.notes_text_frame.text = (
    "This is the most memorable slide — a real detective beat. The brand-name "
    "collision (watches + perfume) is genuinely funny and shows rigour. "
    "Examiners love seeing you question your own data and fix it.")

# ============================================================ 6. TWO CROWDS
s = slide(LIGHT)
text(s, "Two crowds, two very different moods", 0.7, 0.55, 12, 0.9, 36, INK,
     True, HEAD)
text(s, "Casual viewers are upbeat. The app's own users? Three times angrier.",
     0.7, 1.45, 12, 0.5, 19, "1F9D63", True, BODY)
cd = CategoryChartData()
cd.categories = ["Positive", "Neutral", "Negative"]
cd.add_series("YouTube", (61, 29, 10))
cd.add_series("Google Play", (57, 15, 29))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(2.15),
                        Inches(7.4), Inches(4.6), cd)
chart = gf.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(13); chart.legend.font.name = BODY
chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = rgb(YT_C)
chart.series[1].format.fill.solid(); chart.series[1].format.fill.fore_color.rgb = rgb(GP_C)
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
plot.data_labels.font.size = Pt(11); plot.data_labels.font.bold = True
plot.data_labels.font.color.rgb = rgb(INK)
va = chart.value_axis
va.has_major_gridlines = False; va.maximum_scale = 70; va.minimum_scale = 0
va.visible = False
chart.category_axis.tick_labels.font.size = Pt(13)
chart.category_axis.tick_labels.font.bold = True
plats = [("▶ YouTube", YT_C, "10%", "negative · 541 comments"),
         ("▷ Google Play", GP_C, "29%", "negative · 157 reviews")]
px, pw = 8.45, 4.18
for i, (name, col, big2, sub) in enumerate(plats):
    py = 2.15 + i * 2.35
    card(s, px, py, pw, 2.1)
    text(s, name, px + 0.4, py + 0.28, pw - 0.8, 0.5, 18, col, True, HEAD)
    text(s, big2, px + 0.4, py + 0.72, pw - 0.8, 1.0, 44, INK, True, HEAD)
    text(s, sub, px + 0.4, py + 1.6, pw - 0.8, 0.4, 14, MUTED, False, BODY)
s.notes_slide.notes_text_frame.text = (
    "The reveal: where people complain matters. YouTube viewers are watching "
    "videos and mostly upbeat (10% negative). The people who actually downloaded "
    "the app are 3x angrier (29% negative). That sets up the next slide.")

# ============================================================ 7. SMOKING GUN
s = slide(LIGHT)
text(s, "One thing they're shouting about: the app", 0.7, 0.55, 12.4, 0.9, 36,
     INK, True, HEAD)
text(s, "Nearly HALF of all app-store complaints are about the app itself — "
        "crashing, freezing, losing orders.", 0.7, 1.45, 12, 0.6, 18, MUTED,
     False, BODY)
issues = [
    ("▷ Google Play", GP_C, 45,
     [("App & website", 22), ("Orders & collection", 9),
      ("Customer service", 7), ("Stock & availability", 4)]),
    ("▶ YouTube", YT_C, 55,
     [("Orders & collection", 10), ("Customer service", 6),
      ("Delivery & dispatch", 5), ("App & website", 4)]),
]
for i, (name, col, negn, rows) in enumerate(issues):
    x = 0.7 + i * (5.86 + 0.2)
    card(s, x, 2.15, 5.86, 3.45)
    text(s, name, x + 0.45, 2.42, 4.0, 0.5, 19, col, True, HEAD)
    text(s, f"{negn} negative comments", x + 0.45, 2.86, 4.5, 0.4, 12, MUTED,
         False, BODY)
    mx = rows[0][1]
    for j, (theme, cnt) in enumerate(rows):
        ry = 3.4 + j * 0.55
        text(s, theme, x + 0.45, ry, 3.7, 0.35, 14, INK, True, BODY)
        text(s, str(cnt), x + 4.95, ry, 0.55, 0.35, 14, MUTED, True, BODY,
             align=PP_ALIGN.RIGHT)
        bar(s, x + 0.45, ry + 0.34, 4.95, cnt / mx, NEG)
hl = card(s, 0.7, 5.85, 11.93, 1.0, fill=NAVY, shadow=True)
text(s, "The recommendation writes itself:  fix the app first — its own users "
        "are begging for it.", 1.0, 5.85, 11.3, 1.0, 19, WHITE, True, BODY,
     anchor=MSO_ANCHOR.MIDDLE, italic=True)
s.notes_slide.notes_text_frame.text = (
    "The payoff of the analysis: a concrete recommendation. App & website is "
    "22 of 45 Google-Play complaints — nearly half. The people who rely on the "
    "app most are the ones it's failing. That's the headline for Argos.")

# ============================================================ 8. INTEGRITY
s = slide(NAVY)
text(s, "I could have faked it", 0.7, 0.8, 12, 0.9, 38, WHITE, True, HEAD)
text(s, "I wanted X (Twitter) in here too — but reading tweets now costs about "
        "£80 a month, and the free tools are broken.", 0.7, 1.85, 11.9, 0.9,
     20, CYAN, False, BODY, italic=True)
pts = [
    "I could have filled the gap with invented “sample” data.",
    "It would have looked more complete and more impressive.",
    "I didn't — because made-up numbers aren't insight, they're decoration.",
]
yy = 3.0
for p in pts:
    dot(s, 0.8, yy + 0.12, 0.18, CYAN)
    text(s, p, 1.2, yy, 11.0, 0.6, 18, "D7E8EE", False, BODY)
    yy += 0.72
band = card(s, 0.7, 5.5, 11.93, 1.3, fill=TEAL, shadow=True)
text(s, "A good data project shows real data — or none at all.",
     1.0, 5.5, 11.3, 1.3, 22, WHITE, True, BODY, anchor=MSO_ANCHOR.MIDDLE,
     italic=True)
s.notes_slide.notes_text_frame.text = (
    "Counter-intuitively, admitting a limitation builds trust. Explain you "
    "refused to fake X data even though it would have looked better. Integrity "
    "over appearance — this lands really well with examiners.")

# ============================================================ 9. PAYOFF
s = slide(LIGHT)
text(s, "And it never sleeps", 0.7, 0.6, 12, 0.9, 38, INK, True, HEAD)
text(s, "The bot runs every hour in the cloud — my laptop can be switched off. "
        "It refreshes a live dashboard anyone can open… and you can literally "
        "ask it questions in plain English.", 0.7, 1.6, 11.9, 1.1, 19, MUTED,
     False, BODY, line_spacing=1.2)
stats = [("698", "real comments analyzed"), ("£0", "total running cost"),
         ("Every hour", "updates on its own"), ("Ask it", "anything, live")]
gw, gh = 2.93, 2.5
gx = 0.7
for i, (big2, lab) in enumerate(stats):
    x = gx + i * (gw + 0.18)
    card(s, x, 3.1, gw, gh)
    text(s, big2, x, 3.5, gw, 1.1, 40, TEAL, True, HEAD, align=PP_ALIGN.CENTER)
    text(s, lab, x, 4.75, gw, 0.8, 15, MUTED, False, BODY, align=PP_ALIGN.CENTER)
card(s, 0.7, 5.85, 11.93, 0.82, fill=NAVY, shadow=True)
text(s, "🌐  See it live:   https://thefootballbuzz01-creator.github.io/sentiment-analysis-bot/",
     0.9, 5.85, 11.5, 0.82, 16, CYAN, True, BODY, anchor=MSO_ANCHOR.MIDDLE,
     align=PP_ALIGN.CENTER)
s.notes_slide.notes_text_frame.text = (
    "The triumphant beat. It's not a one-off script — it's a living system: "
    "hourly, free, hands-off, and interactive (you can demo the chatbox live "
    "here if you have it open). End on the 'one person, for nothing' line.")

# ============================================================ 10. CLOSING
s = slide(NAVY)
text(s, "The voices were always there.", 0.9, 2.0, 11.5, 1.0, 40, WHITE, True,
     HEAD)
text(s, "We just had to build something that listens.", 0.9, 3.0, 11.5, 1.0,
     34, CYAN, True, HEAD, italic=True)
text(s, "Real data.    Zero cost.    Listening 24/7.", 0.9, 4.6, 11.5, 0.6,
     20, ICE, False, BODY)
text(s, "See it yourself:   https://thefootballbuzz01-creator.github.io/sentiment-analysis-bot/",
     0.9, 5.5, 11.5, 0.5, 16, CYAN, True, BODY)
text(s, "Thank you", 0.9, 6.45, 11.5, 0.5, 16, "9FC6D3", False, BODY)
s.notes_slide.notes_text_frame.text = (
    "Land the message, not the tech. The project was never really about code — "
    "it was about listening to people at a scale humans can't. Deliver the two "
    "lines slowly, then invite questions.")

out = "Sentiment_Analysis_Presentation.pptx"
try:
    prs.save(out)
except PermissionError:
    out = "Sentiment_Analysis_Presentation_UPDATED.pptx"
    prs.save(out)
    print("(original was open/locked — wrote a new copy instead)")
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
